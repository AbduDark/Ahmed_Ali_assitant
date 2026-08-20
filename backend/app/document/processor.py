"""Document processing orchestrator."""

from __future__ import annotations

import os
from pathlib import Path

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.ai.router import get_embedding_provider
from app.config import settings
from app.core.logging import get_logger
from app.document.language_detector import detect_language
from app.document.pdf_extractor import PDFExtractor
from app.document.text_cleaner import TextCleaner
from app.models.chunk import ReferenceChunk
from app.models.reference import Reference, ReferenceStatus
from app.rag.chunker import TextChunker
from app.rag.embedder import EmbeddingService

logger = get_logger(__name__)


class DocumentProcessor:
    """
    Full document processing pipeline:
    Upload → Validate → Extract → Clean → Chunk → Embed → Store
    """

    def __init__(self):
        self.pdf_extractor = PDFExtractor()
        self.text_cleaner = TextCleaner()
        self.chunker = TextChunker(
            chunk_size=settings.rag_chunk_size,
            chunk_overlap=settings.rag_chunk_overlap,
        )

    async def process_reference(
        self,
        reference_id: str,
        db: AsyncSession,
    ) -> None:
        """
        Process a reference document end-to-end.

        This is meant to run as a background job.
        """
        # Load the reference
        result = await db.execute(
            select(Reference).where(Reference.id == reference_id)
        )
        reference = result.scalar_one_or_none()

        if not reference:
            logger.error(f"Reference {reference_id} not found")
            return

        logger.info(f"Processing reference: {reference.title} ({reference.id})")

        try:
            # Mark as processing
            reference.status = ReferenceStatus.PROCESSING
            await db.commit()

            # Step 1: Extract text
            text_pages = await self._extract_text(reference)

            if not text_pages:
                reference.status = ReferenceStatus.FAILED
                reference.error_message = "لم يتم العثور على نص في المستند"
                await db.commit()
                return

            # Step 2: Clean text
            cleaned_pages = [
                (page_num, self.text_cleaner.clean(text))
                for page_num, text in text_pages
            ]

            # Step 3: Detect language
            full_text = " ".join(text for _, text in cleaned_pages[:5])
            language = detect_language(full_text)
            reference.language = language

            # Step 4: Chunk
            chunks_data = self.chunker.chunk_pages(cleaned_pages)
            logger.info(f"Created {len(chunks_data)} chunks from {reference.title}")

            # Step 5: Generate embeddings
            embedding_service = EmbeddingService(get_embedding_provider())
            texts_to_embed = [
                self.text_cleaner.clean_for_embedding(chunk["content"])
                for chunk in chunks_data
            ]
            embeddings = await embedding_service.embed_batch(texts_to_embed)

            # Step 6: Delete old chunks (for reprocessing)
            old_chunks = await db.execute(
                select(ReferenceChunk).where(
                    ReferenceChunk.reference_id == reference_id
                )
            )
            for old_chunk in old_chunks.scalars():
                await db.delete(old_chunk)

            # Step 7: Store new chunks
            for i, chunk_data in enumerate(chunks_data):
                chunk = ReferenceChunk(
                    reference_id=reference_id,
                    content=chunk_data["content"],
                    embedding=embeddings[i] if i < len(embeddings) else None,
                    page_number=chunk_data.get("page_number"),
                    section=chunk_data.get("section"),
                    chunk_index=i,
                    metadata_json={
                        "reference_title": reference.title,
                        "subject_id": reference.subject_id,
                        "unit_id": reference.unit_id,
                        "lesson_id": reference.lesson_id,
                        "language": language,
                    },
                )
                db.add(chunk)

            # Update reference status
            reference.status = ReferenceStatus.READY
            reference.chunk_count = len(chunks_data)
            reference.page_count = len(text_pages)
            reference.error_message = None

            await db.commit()
            logger.info(
                f"Reference {reference.title} processed successfully: "
                f"{len(chunks_data)} chunks, {len(text_pages)} pages"
            )

        except Exception as e:
            logger.error(f"Failed to process reference {reference_id}: {e}")
            reference.status = ReferenceStatus.FAILED
            reference.error_message = str(e)
            await db.commit()
            raise

    async def _extract_text(
        self, reference: Reference
    ) -> list[tuple[int, str]]:
        """Extract text based on file type or source_url. Returns list of (page_num, text)."""
        # If source_url is provided and no local file
        if reference.source_url and not reference.file_path:
            return await self._extract_url(reference.source_url)

        if not reference.file_path:
            return []

        file_path = Path(reference.file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = file_path.suffix.lower()

        if file_ext == ".pdf":
            doc = self.pdf_extractor.extract(file_path)
            return [(p.page_number, p.text) for p in doc.pages if p.text.strip()]

        elif file_ext == ".txt":
            text = file_path.read_text(encoding="utf-8")
            return [(1, text)]

        elif file_ext == ".md":
            text = file_path.read_text(encoding="utf-8")
            return [(1, text)]

        elif file_ext == ".docx":
            return await self._extract_docx(file_path)

        elif file_ext == ".pptx":
            return await self._extract_pptx(file_path)

        else:
            raise ValueError(f"Unsupported file type: {file_ext}")

    async def _extract_url(self, url: str) -> list[tuple[int, str]]:
        """Extract clean text content from a web URL / article."""
        import httpx
        from bs4 import BeautifulSoup

        logger.info(f"Fetching web reference from URL: {url}")
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 "
                "(KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
            ),
            "Accept-Language": "ar,en;q=0.9",
        }

        async with httpx.AsyncClient(timeout=25.0, follow_redirects=True) as client:
            resp = await client.get(url, headers=headers)
            resp.raise_for_status()
            html = resp.text

        soup = BeautifulSoup(html, "html.parser")

        # Remove irrelevant and non-content elements
        for tag in soup(["script", "style", "nav", "footer", "header", "aside", "noscript", "svg", "button", "iframe", "form"]):
            tag.decompose()

        # Try to find main article or content area
        main_content = (
            soup.find("article")
            or soup.find("main")
            or soup.find("div", class_=lambda c: c and any(k in c.lower() for k in ("content", "article", "post", "entry", "body")))
            or soup.body
            or soup
        )

        text = main_content.get_text(separator="\n")
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        cleaned = "\n".join(lines)

        if not cleaned:
            raise ValueError(f"Could not extract meaningful text from URL: {url}")

        logger.info(f"Successfully extracted {len(cleaned)} characters from URL: {url}")
        return [(1, cleaned)]

    async def _extract_docx(self, file_path: Path) -> list[tuple[int, str]]:
        """Extract text from DOCX files."""
        from docx import Document as DocxDocument

        doc = DocxDocument(str(file_path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        full_text = "\n\n".join(paragraphs)
        return [(1, full_text)]

    async def _extract_pptx(self, file_path: Path) -> list[tuple[int, str]]:
        """Extract text from PPTX files."""
        from pptx import Presentation

        prs = Presentation(str(file_path))
        pages: list[tuple[int, str]] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
            if texts:
                pages.append((slide_num, "\n".join(texts)))

        return pages
